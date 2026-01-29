"""
文档处理服务
用于提取各种文档格式的文本内容
"""
import os
from pathlib import Path
from typing import List, Dict, Any
import PyPDF2
from docx import Document
import json


class DocumentProcessor:
    """文档处理器"""
    
    def __init__(self):
        self.chunk_size = 1000  # 每个文本块的字符数
        self.chunk_overlap = 200  # 文本块之间的重叠字符数
    
    def extract_text(self, file_path: str, file_type: str) -> str:
        """
        从文档中提取文本
        
        Args:
            file_path: 文件路径
            file_type: 文件类型扩展名 (如 .pdf, .docx)
        
        Returns:
            提取的文本内容
        """
        try:
            if file_type.lower() == '.pdf':
                return self._extract_from_pdf(file_path)
            elif file_type.lower() in ['.doc', '.docx']:
                return self._extract_from_word(file_path)
            elif file_type.lower() in ['.ppt', '.pptx']:
                return self._extract_from_ppt(file_path)
            elif file_type.lower() in ['.txt', '.md']:
                return self._extract_from_text(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {file_type}")
        except Exception as e:
            raise Exception(f"文本提取失败: {str(e)}")
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """从PDF提取文本"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n\n=== 第 {page_num + 1} 页 ===\n{page_text}"
        except Exception as e:
            raise Exception(f"PDF解析失败: {str(e)}")
        
        return text.strip()
    
    def _extract_from_word(self, file_path: str) -> str:
        """从Word文档提取文本"""
        try:
            doc = Document(file_path)
            paragraphs = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        paragraphs.append(row_text)
            
            return '\n\n'.join(paragraphs)
        except Exception as e:
            raise Exception(f"Word文档解析失败: {str(e)}")
    
    def _extract_from_ppt(self, file_path: str) -> str:
        """从PowerPoint提取文本"""
        try:
            # PPT文件需要python-pptx库
            # 如果未安装，返回提示信息
            try:
                from pptx import Presentation
            except ImportError:
                return "[PowerPoint文件内容提取需要安装python-pptx库]"
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n\n=== 幻灯片 {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    text_parts.append(slide_text)
            
            return '\n'.join(text_parts)
        except ImportError:
            return "[PowerPoint文件内容提取需要安装python-pptx库]"
        except Exception as e:
            raise Exception(f"PowerPoint解析失败: {str(e)}")
    
    def _extract_from_text(self, file_path: str) -> str:
        """从文本文件提取内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # 如果UTF-8失败，尝试GBK编码（中文Windows常用）
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    return file.read()
            except Exception as e:
                raise Exception(f"文本文件读取失败: {str(e)}")
    
    def split_text_into_chunks(self, text: str, metadata: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """
        将长文本切分成多个文本块
        
        Args:
            text: 要切分的文本
            metadata: 附加的元数据
        
        Returns:
            文本块列表，每个块包含text和metadata
        """
        if not text or not text.strip():
            return []
        
        chunks = []
        text = text.strip()
        
        # 如果文本较短，直接返回一个块
        if len(text) <= self.chunk_size:
            chunks.append({
                'text': text,
                'chunk_index': 0,
                'metadata': metadata or {}
            })
            return chunks
        
        # 切分长文本
        start = 0
        chunk_index = 0
        
        while start < len(text):
            # 确定结束位置
            end = start + self.chunk_size
            
            # 如果不是最后一块，尝试在句子边界处切分
            if end < len(text):
                # 查找最近的句子结束标记
                sentence_endings = ['。', '！', '？', '.\n', '!\n', '?\n', '\n\n']
                best_end = end
                
                for i in range(end, start + self.chunk_size // 2, -1):
                    for ending in sentence_endings:
                        if text[i:i+len(ending)] == ending:
                            best_end = i + len(ending)
                            break
                    if best_end != end:
                        break
                
                end = best_end
            
            # 提取文本块
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunk_metadata = (metadata or {}).copy()
                chunk_metadata.update({
                    'chunk_index': chunk_index,
                    'start_pos': start,
                    'end_pos': end,
                    'total_chunks': 'calculating'  # 稍后更新
                })
                
                chunks.append({
                    'text': chunk_text,
                    'chunk_index': chunk_index,
                    'metadata': chunk_metadata
                })
                
                chunk_index += 1
            
            # 移动到下一块，考虑重叠
            start = end - self.chunk_overlap
            if start < 0:
                start = 0
        
        # 更新总块数
        total_chunks = len(chunks)
        for chunk in chunks:
            chunk['metadata']['total_chunks'] = total_chunks
        
        return chunks
    
    def process_document(self, file_path: str, file_type: str, document_metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        处理文档：提取文本并切分成块
        
        Args:
            file_path: 文件路径
            file_type: 文件类型
            document_metadata: 文档元数据
        
        Returns:
            包含原始文本和文本块的字典
        """
        # 提取文本
        text = self.extract_text(file_path, file_type)
        
        if not text or not text.strip():
            raise Exception("文档内容为空或无法提取")
        
        # 切分文本
        chunks = self.split_text_into_chunks(text, document_metadata)
        
        return {
            'full_text': text,
            'chunks': chunks,
            'total_chars': len(text),
            'total_chunks': len(chunks)
        }


# 创建全局实例
document_processor = DocumentProcessor()
